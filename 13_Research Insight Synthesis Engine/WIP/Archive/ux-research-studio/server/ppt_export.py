from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import Dict

from pptx import Presentation

from .config import STUDIO_DIR
from . import interviews as interview_assets
from .modules_meta import MODULES
from .workspace import read_research_file, research_dir, safe_project_path

TEMPLATE_PATH = STUDIO_DIR / "templates" / "UX_Research_Report_TEMPLATE.pptx"
MAX_FIELD_CHARS = 3500


def _truncate(text: str, limit: int = MAX_FIELD_CHARS) -> str:
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return text[: limit - 20] + "\n\n…（已截断）"


def _read_module_md(project_relative: str, filename: str) -> str:
    try:
        return read_research_file(project_relative, filename)
    except (FileNotFoundError, ValueError):
        return ""


def _strip_md(md: str) -> str:
    if not md:
        return ""
    t = re.sub(r"```[\s\S]*?```", "", md)
    t = re.sub(r"^#+\s*", "", t, flags=re.MULTILINE)
    t = re.sub(r"\*\*([^*]+)\*\*", r"\1", t)
    return t.strip()


def build_placeholder_map(project_relative: str, context: dict) -> Dict[str, str]:
    """从项目 Research 文件与四问构建占位符映射。"""
    title = context.get("topic") or project_relative.split("/")[-1] or "UX Research"
    mapping = {
        "{{PROJECT_TITLE}}": title,
        "{{SUBTITLE}}": "UX Research Planning · 研究规划汇报",
        "{{TOPIC}}": context.get("topic") or "（待填）",
        "{{AUDIENCE}}": context.get("audience") or "（待填）",
        "{{STAGE}}": context.get("stage") or "（待填）",
        "{{CORE_QUESTION}}": context.get("core_question") or "（待填）",
        "{{EXPORT_TIME}}": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "{{PROGRESS}}": "—",
    }

    for meta in MODULES:
        key = f"{{{{MODULE_{meta.id}}}}}"
        if meta.output_is_dir:
            combined = interview_assets.combined_transcripts(project_relative)
            body = combined or _read_module_md(project_relative, f"{meta.output_file}/reference_01.md")
        else:
            body = _read_module_md(project_relative, meta.output_file)
        mapping[key] = _truncate(_strip_md(body)) or "（本模块暂无内容，请先在 Studio 生成并保存）"

    try:
        from .workspace import progress_summary

        ps = progress_summary(project_relative)
        done = ps.get("completed") or []
        mapping["{{PROGRESS}}"] = (
            f"{ps.get('completed_count', 0)}/9 · 模块 {', '.join(str(x) for x in done) or '无'}"
        )
    except ValueError:
        pass

    return mapping


def _replace_in_shape(shape, mapping: Dict[str, str]) -> None:
    if not shape.has_text_frame:
        return
    text = shape.text_frame.text or ""
    if "{{" not in text:
        return
    new = text
    for k, v in mapping.items():
        new = new.replace(k, v)
    if new == text:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = new


def export_ppt(project_relative: str, context: dict) -> Path:
    if not TEMPLATE_PATH.is_file():
        raise FileNotFoundError(
            f"模版不存在：{TEMPLATE_PATH}，请运行 scripts/build_ppt_template.py"
        )
    safe_project_path(project_relative)
    mapping = build_placeholder_map(project_relative, context)

    prs = Presentation(str(TEMPLATE_PATH))
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_in_shape(shape, mapping)

    out_dir = research_dir(project_relative) / "exports"
    out_dir.mkdir(parents=True, exist_ok=True)
    safe_name = re.sub(r"[^\w\u4e00-\u9fff\-]+", "_", mapping["{{PROJECT_TITLE}}"])[:40]
    out_path = out_dir / f"UX_Research_Report_{safe_name}.pptx"
    prs.save(str(out_path))
    return out_path
