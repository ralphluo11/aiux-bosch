#!/usr/bin/env python3
"""生成 UX Research 汇报用假模版（占位符 {{KEY}}）。"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "templates" / "UX_Research_Report_TEMPLATE.pptx"
SKILL_OUT = (
    Path(__file__).resolve().parent.parent.parent
    / ".cursor"
    / "skills"
    / "ux-research-planning"
    / "templates"
    / "UX_Research_Report_TEMPLATE.pptx"
)


def _add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_content_slide(prs, title, body_placeholder):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if body is not None:
        body.text = body_placeholder


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "{{PROJECT_TITLE}}",
        "{{SUBTITLE}}\n课题：{{TOPIC}} · 用户：{{AUDIENCE}}",
    )
    _add_content_slide(
        prs,
        "研究目标与核心问题",
        "{{MODULE_1}}\n\n核心问题：{{CORE_QUESTION}}",
    )
    _add_content_slide(prs, "桌面研究 · 利益相关者", "{{MODULE_2}}\n\n---\n\n{{MODULE_3}}")
    _add_content_slide(prs, "招募与访纲", "{{MODULE_4}}\n\n---\n\n{{MODULE_5}}")
    _add_content_slide(prs, "素材与画像", "{{MODULE_7}}")
    _add_content_slide(
        prs,
        "旅程与机会点",
        "{{MODULE_8}}",
    )
    _add_content_slide(prs, "概念方向", "{{MODULE_9}}")
    _add_content_slide(
        prs,
        "项目进度",
        "已完成模块：{{PROGRESS}}\n\n导出时间：{{EXPORT_TIME}}",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")

    SKILL_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(SKILL_OUT))
    print(f"Wrote {SKILL_OUT}")


if __name__ == "__main__":
    main()
